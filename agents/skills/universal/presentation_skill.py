"""
Universal Skill: Presentation (PPTX)

Generates a PowerPoint slide deck from live project data.
Slides:
  1. Cover (project name, date, produced by)
  2. Executive Summary (health, key metrics)
  3. KPI Scorecard (table of KPIs with status)
  4. Milestone Timeline (table of milestones)
  5. Budget Overview (visual bar)
  6. Risk Summary (open high-impact risks)
  7. Next Steps / Actions Pending

Trigger phrases: "presentation", "slide deck", "pptx", "slides",
                 "deck", "board presentation"
"""

from __future__ import annotations

import logging
import os
import sqlite3
import uuid
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from agents.skills.base_skill import BaseSkill, SkillOutput
from infrastructure.project_context import ProjectContext

logger = logging.getLogger(__name__)

# Brand colours
DARK_BLUE  = RGBColor(0x1F, 0x35, 0x64)
ACCENT     = RGBColor(0x2E, 0x75, 0xB6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x70, 0xAD, 0x47)
RED        = RGBColor(0xC0, 0x00, 0x00)
AMBER      = RGBColor(0xFF, 0xC0, 0x00)


def _status_color(status: str) -> RGBColor:
    s = (status or "").lower()
    if s in ("on_track", "completed", "active"):
        return GREEN
    if s in ("off_track", "overdue", "critical"):
        return RED
    return AMBER


def _add_slide(prs: Presentation, layout_idx: int = 6):
    """Add a blank slide."""
    layout = prs.slide_layouts[layout_idx]   # 6 = blank
    return prs.slides.add_slide(layout)


def _text_box(slide, left, top, width, height, text, font_size=18,
              bold=False, color: RGBColor = None, align=PP_ALIGN.LEFT):
    tb  = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf  = tb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return tb


def _fill_shape(shape, color: RGBColor):
    from pptx.util import Pt
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


class PresentationSkill(BaseSkill):
    skill_id        = "project_presentation"
    dept_id         = "all"
    title_template  = "{project_name} — Executive Presentation"
    description     = "Generate a PPTX slide deck: executive summary, KPIs, milestones, budget, risks, next steps."
    output_format   = "pptx"
    trigger_phrases = [
        "presentation", "slide deck", "slides", "pptx", "deck",
        "board presentation", "executive presentation", "generate slides",
        "create presentation", "make a deck",
    ]

    async def execute(self, context: ProjectContext, params: dict[str, Any] = None) -> SkillOutput:
        params = params or {}
        title  = params.get("title") or self._make_title(context)

        try:
            db_path = getattr(context, "db_path", None)

            # ── Fetch project data ─────────────────────────────────────────
            proj_info   = {}
            kpis        = []
            milestones  = []
            risks       = []
            actions     = []

            if db_path and os.path.exists(db_path):
                conn = sqlite3.connect(f"file:{db_path}?mode=ro", uri=True, timeout=10)
                conn.row_factory = sqlite3.Row
                try:
                    r = conn.execute(
                        "SELECT name, status, health_status, budget_total, budget_spent, "
                        "start_date, target_end_date, description "
                        "FROM project_details LIMIT 1"
                    ).fetchone()
                    if r:
                        proj_info = dict(r)

                    kpis = conn.execute(
                        "SELECT name, current_value, target_value, unit, status "
                        "FROM project_kpis ORDER BY status DESC LIMIT 8"
                    ).fetchall()

                    milestones = conn.execute(
                        "SELECT name, due_date, status, owner "
                        "FROM project_milestones ORDER BY due_date ASC LIMIT 8"
                    ).fetchall()

                    risks = conn.execute(
                        "SELECT title, severity, impact, likelihood, status "
                        "FROM project_risks WHERE status='open' "
                        "ORDER BY severity DESC LIMIT 6"
                    ).fetchall()

                    actions = conn.execute(
                        "SELECT title, priority, status FROM agent_action_queue "
                        "WHERE status='pending' ORDER BY priority DESC LIMIT 5"
                    ).fetchall()
                finally:
                    conn.close()

            # ── Build presentation ─────────────────────────────────────────
            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            now       = datetime.utcnow().strftime("%B %d, %Y")
            proj_name = proj_info.get("name") or context.project_id[:12]
            health    = proj_info.get("health_status") or "unknown"

            # ── Slide 1: Cover ─────────────────────────────────────────────
            slide = _add_slide(prs)
            bg = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                Inches(0), Inches(0), Inches(13.33), Inches(7.5)
            )
            _fill_shape(bg, DARK_BLUE)
            bg.line.fill.background()

            _text_box(slide, 0.5, 2.2, 12, 1.2, proj_name,
                      font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            _text_box(slide, 0.5, 3.6, 12, 0.6, title,
                      font_size=22, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
            _text_box(slide, 0.5, 4.4, 12, 0.5, f"Generated: {now}  |  Status: {health.upper()}",
                      font_size=14, color=LIGHT_GREY, align=PP_ALIGN.CENTER)

            # ── Slide 2: Executive Summary ─────────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, "Executive Summary",
                      font_size=28, bold=True, color=DARK_BLUE)
            _add_divider(slide)

            budget_total = proj_info.get("budget_total") or 0
            budget_spent = proj_info.get("budget_spent") or 0
            budget_pct   = (budget_spent / budget_total * 100) if budget_total else 0

            summary_lines = [
                f"Project Health:  {health.upper()}",
                f"Budget Used:     ${budget_spent:,.0f} of ${budget_total:,.0f}  ({budget_pct:.1f}%)",
                f"Start Date:      {proj_info.get('start_date') or 'TBD'}",
                f"Target End:      {proj_info.get('target_end_date') or 'TBD'}",
                f"KPIs Tracked:    {len(kpis)}",
                f"Open Risks:      {len(risks)}",
                f"Pending Actions: {len(actions)}",
            ]
            body = "\n".join(summary_lines)
            _text_box(slide, 0.4, 1.0, 12, 5.5, body, font_size=18, color=DARK_BLUE)

            # ── Slide 3: KPI Scorecard ─────────────────────────────────────
            if kpis:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "KPI Scorecard",
                          font_size=28, bold=True, color=DARK_BLUE)
                _add_divider(slide)

                rows_data = [["KPI", "Current", "Target", "Unit", "Status"]]
                for k in kpis:
                    rows_data.append([
                        k["name"], str(k["current_value"] or "—"),
                        str(k["target_value"] or "—"), str(k["unit"] or ""),
                        str(k["status"] or "")
                    ])
                _add_table(slide, rows_data, left=0.4, top=1.0, width=12.5, row_height=0.45)

            # ── Slide 4: Milestones ────────────────────────────────────────
            if milestones:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "Milestone Timeline",
                          font_size=28, bold=True, color=DARK_BLUE)
                _add_divider(slide)

                rows_data = [["Milestone", "Due Date", "Owner", "Status"]]
                for m in milestones:
                    rows_data.append([
                        m["name"], str(m["due_date"] or "TBD"),
                        str(m["owner"] or "Unassigned"), str(m["status"] or "")
                    ])
                _add_table(slide, rows_data, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # ── Slide 5: Risk Summary ──────────────────────────────────────
            if risks:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "Risk Summary",
                          font_size=28, bold=True, color=DARK_BLUE)
                _add_divider(slide)

                rows_data = [["Risk", "Severity", "Impact", "Likelihood", "Status"]]
                for r in risks:
                    rows_data.append([
                        r["title"], str(r["severity"] or ""),
                        str(r["impact"] or ""), str(r["likelihood"] or ""),
                        str(r["status"] or "")
                    ])
                _add_table(slide, rows_data, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # ── Slide 6: Next Steps ────────────────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, "Pending Actions & Next Steps",
                      font_size=28, bold=True, color=DARK_BLUE)
            _add_divider(slide)

            if actions:
                rows_data = [["Action", "Priority", "Status"]]
                for a in actions:
                    rows_data.append([a["title"], str(a["priority"] or ""), str(a["status"] or "")])
                _add_table(slide, rows_data, left=0.4, top=1.0, width=10, row_height=0.5)
            else:
                _text_box(slide, 0.4, 1.2, 12, 1, "No pending actions at this time.",
                          font_size=18, color=DARK_BLUE)

            # ── Save ───────────────────────────────────────────────────────
            out_dir  = self._output_dir(context)
            filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
            out_path = os.path.join(out_dir, filename)
            prs.save(out_path)

            n_slides = len(prs.slides)
            return SkillOutput(
                skill_id    = self.skill_id,
                dept_id     = self.dept_id,
                title       = title,
                file_format = self.output_format,
                file_path   = out_path,
                preview     = f"Presentation: {n_slides} slides covering KPIs, milestones, risks, and next steps.",
                pages       = n_slides,
            )

        except Exception as e:
            logger.error(f"[PresentationSkill] Failed: {e}", exc_info=True)
            return SkillOutput(
                skill_id    = self.skill_id,
                dept_id     = self.dept_id,
                title       = title,
                file_format = self.output_format,
                error       = str(e),
            )


# ── Helpers ───────────────────────────────────────────────────────────────────

def _add_divider(slide, top: float = 0.85):
    """Thin horizontal line under the slide title."""
    from pptx.util import Inches, Pt, Emu
    line = slide.shapes.add_shape(1, Inches(0.4), Inches(top), Inches(12.5), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()


def _add_table(slide, rows_data, left=0.4, top=1.0, width=12.5, row_height=0.45):
    """Add a styled table to a slide."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0]) if rows_data else 1

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(left), Inches(top),
        Inches(width), Inches(row_height * n_rows),
    )
    tbl = table_shape.table

    for r_idx, row_vals in enumerate(rows_data):
        for c_idx, val in enumerate(row_vals):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(val)
            p   = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(13)
            if r_idx == 0:
                run.font.bold  = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
            else:
                run.font.color.rgb = DARK_BLUE
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_GREY


# ── Self-registration ─────────────────────────────────────────────────────────

def _register(registry) -> None:
    registry.register(PresentationSkill())
