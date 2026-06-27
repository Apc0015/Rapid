"""
Finance Skill: Board Financial Presentation (PPTX)

Generates an executive-ready PPTX for board or leadership meetings:
  - Cover slide
  - Financial Highlights
  - Budget vs. Actual (all categories)
  - KPI Performance
  - Risk & Issues
  - Recommendations

Trigger phrases: "board presentation", "board deck", "executive presentation",
                 "financial presentation", "leadership deck"
"""

from __future__ import annotations

import logging
import os
import sqlite3
import uuid
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from agents.skills.base_skill import BaseSkill, SkillOutput
from infrastructure.project_context import ProjectContext
# Re-use helper functions from presentation_skill
from agents.skills.universal.presentation_skill import (
    _add_slide, _text_box, _add_divider, _add_table,
    DARK_BLUE, ACCENT, WHITE, LIGHT_GREY, GREEN, RED, AMBER,
)

logger = logging.getLogger(__name__)


class BoardPresentationSkill(BaseSkill):
    skill_id        = "board_presentation"
    dept_id         = "finance"
    title_template  = "{project_name} — Board Financial Presentation"
    description     = "Finance: Generate an executive PPTX for board/leadership: financial highlights, budget vs. actual, KPIs, risks, recommendations."
    output_format   = "pptx"
    trigger_phrases = [
        "board presentation", "board deck", "executive presentation",
        "financial presentation", "leadership deck", "board report",
        "executive deck", "board financial", "finance board",
    ]

    async def execute(self, context: ProjectContext, params: dict[str, Any] = None) -> SkillOutput:
        params = params or {}
        title  = params.get("title") or self._make_title(context)

        try:
            db_path   = getattr(context, "db_path", None)
            proj_info = {}
            kpis      = []
            budget    = []
            risks     = []

            if db_path and os.path.exists(db_path):
                conn = sqlite3.connect(f"file:{db_path}?mode=ro", uri=True, timeout=10)
                conn.row_factory = sqlite3.Row
                try:
                    r = conn.execute("SELECT * FROM project_details LIMIT 1").fetchone()
                    if r:
                        proj_info = dict(r)

                    kpis = [dict(r) for r in conn.execute(
                        "SELECT name, current_value, target_value, unit, status "
                        "FROM project_kpis ORDER BY status DESC LIMIT 8"
                    ).fetchall()]

                    budget = [dict(r) for r in conn.execute(
                        "SELECT category, allocated, spent FROM project_budget_lines ORDER BY category"
                    ).fetchall()]

                    risks = [dict(r) for r in conn.execute(
                        "SELECT title, severity, impact, status FROM project_risks "
                        "WHERE status='open' ORDER BY severity DESC LIMIT 6"
                    ).fetchall()]
                finally:
                    conn.close()

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            now       = datetime.utcnow().strftime("%B %d, %Y")
            proj_name = proj_info.get("name") or context.project_id[:12]
            health    = (proj_info.get("health_status") or "unknown").upper()

            total_alloc = sum(b.get("allocated") or 0 for b in budget)
            total_spent = sum(b.get("spent") or 0 for b in budget)
            total_remain = total_alloc - total_spent
            budget_pct  = (total_spent / total_alloc * 100) if total_alloc else 0

            from pptx.dml.color import RGBColor as _RGB
            from pptx.util import Inches as _I

            # ── Slide 1: Cover ─────────────────────────────────────────────
            slide = _add_slide(prs)
            bg = slide.shapes.add_shape(1, _I(0), _I(0), _I(13.33), _I(7.5))
            bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
            _text_box(slide, 0.5, 1.8, 12, 1.2, proj_name, 40, True, WHITE, 2)
            _text_box(slide, 0.5, 3.2, 12, 0.6, "Board Financial Presentation", 22, False, WHITE, 2)
            _text_box(slide, 0.5, 4.0, 12, 0.5, f"{now}  |  Health: {health}", 14, False, LIGHT_GREY, 2)

            from pptx.enum.text import PP_ALIGN
            # ── Slide 2: Financial Highlights ──────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, "Financial Highlights", 28, True, DARK_BLUE)
            _add_divider(slide)

            summary = (
                f"Total Budget:       ${total_alloc:>12,.2f}\n"
                f"Amount Spent:       ${total_spent:>12,.2f}\n"
                f"Remaining Budget:   ${total_remain:>12,.2f}\n"
                f"Budget Utilization: {budget_pct:>11.1f}%\n"
                f"Health Status:      {health:>12}\n"
                f"Open Risks:         {len(risks):>12}"
            )
            _text_box(slide, 0.4, 1.0, 12, 5.5, summary, 20, False, DARK_BLUE)

            # ── Slide 3: Budget vs. Actual ─────────────────────────────────
            if budget:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "Budget vs. Actual", 28, True, DARK_BLUE)
                _add_divider(slide)

                rows = [["Category", "Allocated ($)", "Spent ($)", "Remaining ($)", "Used %"]]
                for b in budget:
                    alloc = b.get("allocated") or 0
                    spent = b.get("spent") or 0
                    rem   = alloc - spent
                    pct   = f"{spent/alloc*100:.1f}%" if alloc else "N/A"
                    rows.append([b.get("category"), f"${alloc:,.0f}", f"${spent:,.0f}", f"${rem:,.0f}", pct])
                rows.append(["TOTAL", f"${total_alloc:,.0f}", f"${total_spent:,.0f}", f"${total_remain:,.0f}", f"{budget_pct:.1f}%"])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.45)

            # ── Slide 4: KPI Performance ───────────────────────────────────
            if kpis:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "KPI Performance", 28, True, DARK_BLUE)
                _add_divider(slide)

                rows = [["KPI", "Actual", "Target", "Unit", "Status"]]
                for k in kpis:
                    rows.append([k.get("name"), str(k.get("current_value") or "—"),
                                 str(k.get("target_value") or "—"), str(k.get("unit") or ""),
                                 str(k.get("status") or "")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # ── Slide 5: Risk Summary ──────────────────────────────────────
            if risks:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "Risk Summary", 28, True, DARK_BLUE)
                _add_divider(slide)

                rows = [["Risk", "Severity", "Impact", "Status"]]
                for r in risks:
                    rows.append([r.get("title"), str(r.get("severity") or ""),
                                 str(r.get("impact") or ""), str(r.get("status") or "")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # ── Slide 6: Recommendations ───────────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, "Recommendations", 28, True, DARK_BLUE)
            _add_divider(slide)

            recs = []
            if budget_pct >= 90:
                recs.append("🔴 Budget critically consumed — immediate reallocation or scope review needed")
            elif budget_pct >= 80:
                recs.append("🟡 Budget utilization high — monitor spend closely and review forecasts")
            else:
                recs.append("🟢 Budget on track — continue with planned spend schedule")

            high_risks = [r for r in risks if str(r.get("severity") or "").lower() in ("high","critical")]
            if high_risks:
                recs.append(f"⚠️ {len(high_risks)} high/critical risk(s) require immediate mitigation planning")

            on_track_kpis  = sum(1 for k in kpis if str(k.get("status") or "").lower() == "on_track")
            off_track_kpis = sum(1 for k in kpis if str(k.get("status") or "").lower() == "off_track")
            if off_track_kpis:
                recs.append(f"📊 {off_track_kpis} KPI(s) are off track — escalate to department lead")
            if on_track_kpis == len(kpis) and kpis:
                recs.append("✅ All KPIs on track — strong performance this period")

            recs.append("📋 All agent-generated actions await human review and approval in RAPID")

            _text_box(slide, 0.4, 1.0, 12, 5.5, "\n\n".join(recs), 18, False, DARK_BLUE)

            # Save
            out_dir  = self._output_dir(context)
            filename = f"board_presentation_{uuid.uuid4().hex[:8]}.pptx"
            out_path = os.path.join(out_dir, filename)
            prs.save(out_path)

            n_slides = len(prs.slides)
            return SkillOutput(
                skill_id    = self.skill_id,
                dept_id     = self.dept_id,
                title       = title,
                file_format = self.output_format,
                file_path   = out_path,
                preview     = f"Board presentation: {n_slides} slides. Budget {budget_pct:.1f}% used.",
                pages       = n_slides,
            )

        except Exception as e:
            logger.error(f"[BoardPresentationSkill] Failed: {e}", exc_info=True)
            return SkillOutput(skill_id=self.skill_id, dept_id=self.dept_id,
                               title=title, file_format=self.output_format, error=str(e))


def _register(registry) -> None:
    registry.register(BoardPresentationSkill())
