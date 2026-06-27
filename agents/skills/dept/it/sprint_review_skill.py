"""
IT Skill: Sprint Review (PPTX)

Generates a sprint review presentation:
  - Sprint summary (velocity, completed stories, carry-over)
  - Milestone status (completed vs open)
  - System health KPIs
  - Next sprint plan

Trigger phrases: "sprint review", "sprint report", "sprint retrospective",
                 "sprint summary", "iteration review"
"""

from __future__ import annotations

import logging
import os
import sqlite3
import uuid
from datetime import datetime
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from agents.skills.base_skill import BaseSkill, SkillOutput
from infrastructure.project_context import ProjectContext
from agents.skills.universal.presentation_skill import (
    _add_slide, _text_box, _add_divider, _add_table,
    DARK_BLUE, ACCENT, WHITE, LIGHT_GREY, GREEN, RED, AMBER,
)
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)


class SprintReviewSkill(BaseSkill):
    skill_id        = "sprint_review"
    dept_id         = "it"
    title_template  = "{project_name} — Sprint Review"
    description     = "IT: Generate a sprint review PPTX with velocity, completed stories, carry-over, and next sprint plan."
    output_format   = "pptx"
    trigger_phrases = [
        "sprint review", "sprint report", "sprint retrospective",
        "sprint summary", "iteration review", "agile review",
        "sprint deck", "end of sprint",
    ]

    async def execute(self, context: ProjectContext, params: dict[str, Any] = None) -> SkillOutput:
        params = params or {}
        title  = params.get("title") or self._make_title(context)

        try:
            db_path   = getattr(context, "db_path", None)
            proj_info = {}
            kpis      = []
            milestones = []
            risks     = []

            if db_path and os.path.exists(db_path):
                conn = sqlite3.connect(f"file:{db_path}?mode=ro", uri=True, timeout=10)
                conn.row_factory = sqlite3.Row
                try:
                    r = conn.execute("SELECT * FROM project_details LIMIT 1").fetchone()
                    if r:
                        proj_info = dict(r)

                    kpis = [dict(r) for r in conn.execute(
                        "SELECT name, current_value, target_value, unit, status "
                        "FROM project_kpis ORDER BY status DESC LIMIT 10"
                    ).fetchall()]

                    milestones = [dict(r) for r in conn.execute(
                        "SELECT name, due_date, completed_date, status, owner "
                        "FROM project_milestones ORDER BY due_date ASC LIMIT 12"
                    ).fetchall()]

                    risks = [dict(r) for r in conn.execute(
                        "SELECT title, severity, status FROM project_risks "
                        "WHERE status='open' ORDER BY severity DESC LIMIT 5"
                    ).fetchall()]
                finally:
                    conn.close()

            prs = Presentation()
            prs.slide_width  = Inches(13.33)
            prs.slide_height = Inches(7.5)

            now       = datetime.utcnow().strftime("%B %d, %Y")
            proj_name = proj_info.get("name") or context.project_id[:12]
            health    = (proj_info.get("health_status") or "unknown").upper()

            completed_ms = [m for m in milestones if str(m.get("status") or "").lower() == "completed"]
            open_ms      = [m for m in milestones if str(m.get("status") or "").lower() != "completed"]

            # Sprint KPIs
            sprint_kpis = [k for k in kpis if any(
                w in str(k.get("name","")).lower()
                for w in ["velocity", "story", "point", "sprint", "burndown", "bug", "defect",
                          "deploy", "release", "uptime", "incident", "sla"]
            )]

            # ── Slide 1: Cover ─────────────────────────────────────────────
            slide = _add_slide(prs)
            bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
            bg.fill.solid(); bg.fill.fore_color.rgb = DARK_BLUE; bg.line.fill.background()
            _text_box(slide, 0.5, 2.0, 12, 1.0, proj_name, 36, True, WHITE, PP_ALIGN.CENTER)
            _text_box(slide, 0.5, 3.2, 12, 0.6, "Sprint Review", 24, False, WHITE, PP_ALIGN.CENTER)
            _text_box(slide, 0.5, 4.0, 12, 0.5, f"{now}  |  IT Department  |  Health: {health}",
                      14, False, LIGHT_GREY, PP_ALIGN.CENTER)

            # ── Slide 2: Sprint Summary ────────────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, "Sprint Summary", 28, True, DARK_BLUE)
            _add_divider(slide)

            summary = (
                f"Project Health:      {health}\n"
                f"Total Milestones:    {len(milestones)}\n"
                f"Completed This Sprint: {len(completed_ms)}\n"
                f"Carry-over (Open):   {len(open_ms)}\n"
                f"Open Risks:          {len(risks)}\n"
                f"Sprint KPIs Tracked: {len(sprint_kpis)}"
            )
            _text_box(slide, 0.4, 1.0, 12, 5.5, summary, 20, False, DARK_BLUE)

            # ── Slide 3: Completed This Sprint ─────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, f"✅ Completed ({len(completed_ms)})",
                      28, True, DARK_BLUE)
            _add_divider(slide)
            if completed_ms:
                rows = [["Milestone", "Completed Date", "Owner"]]
                for m in completed_ms[:8]:
                    rows.append([m.get("name"), str(m.get("completed_date") or "—"), str(m.get("owner") or "—")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)
            else:
                _text_box(slide, 0.4, 1.2, 12, 1, "No milestones completed this sprint.", 18, False, DARK_BLUE)

            # ── Slide 4: Carry-over ────────────────────────────────────────
            slide = _add_slide(prs)
            _text_box(slide, 0.4, 0.2, 12, 0.6, f"⏭ Carry-over / Next Sprint ({len(open_ms)})",
                      28, True, DARK_BLUE)
            _add_divider(slide)
            if open_ms:
                rows = [["Milestone", "Due Date", "Owner", "Status"]]
                for m in open_ms[:8]:
                    rows.append([m.get("name"), str(m.get("due_date") or "TBD"),
                                 str(m.get("owner") or "—"), str(m.get("status") or "")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)
            else:
                _text_box(slide, 0.4, 1.2, 12, 1, "No carry-over items.", 18, False, DARK_BLUE)

            # ── Slide 5: System Health KPIs ────────────────────────────────
            if sprint_kpis:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "System Health & Velocity", 28, True, DARK_BLUE)
                _add_divider(slide)
                rows = [["KPI", "Current", "Target", "Unit", "Status"]]
                for k in sprint_kpis:
                    rows.append([k.get("name"), str(k.get("current_value") or "—"),
                                 str(k.get("target_value") or "—"), str(k.get("unit") or ""),
                                 str(k.get("status") or "")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # ── Slide 6: Risks ─────────────────────────────────────────────
            if risks:
                slide = _add_slide(prs)
                _text_box(slide, 0.4, 0.2, 12, 0.6, "Open Risks", 28, True, DARK_BLUE)
                _add_divider(slide)
                rows = [["Risk", "Severity", "Status"]]
                for r in risks:
                    rows.append([r.get("title"), str(r.get("severity") or ""), str(r.get("status") or "")])
                _add_table(slide, rows, left=0.4, top=1.0, width=12.5, row_height=0.5)

            # Save
            out_dir  = self._output_dir(context)
            filename = f"sprint_review_{uuid.uuid4().hex[:8]}.pptx"
            out_path = os.path.join(out_dir, filename)
            prs.save(out_path)

            n_slides = len(prs.slides)
            return SkillOutput(
                skill_id    = self.skill_id,
                dept_id     = self.dept_id,
                title       = title,
                file_format = self.output_format,
                file_path   = out_path,
                preview     = f"Sprint review: {len(completed_ms)} completed, {len(open_ms)} carry-over, {n_slides} slides.",
                pages       = n_slides,
            )

        except Exception as e:
            logger.error(f"[SprintReviewSkill] Failed: {e}", exc_info=True)
            return SkillOutput(skill_id=self.skill_id, dept_id=self.dept_id,
                               title=title, file_format=self.output_format, error=str(e))


def _register(registry) -> None:
    registry.register(SprintReviewSkill())
